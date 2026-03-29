"""Embeddings generation and management using OpenAI."""

import io
import logging
from typing import List, Optional

from openai import OpenAI
from sqlalchemy.orm import Session

from mvp.config import settings
from mvp.models import File, FileEmbedding

logger = logging.getLogger(__name__)

# Chunk size for text splitting (in characters)
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200


def extract_text_from_pdf(content: bytes) -> Optional[str]:
    """Extract text from PDF content."""
    try:
        import pdfplumber

        with pdfplumber.open(io.BytesIO(content)) as pdf:
            text_parts = []
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            return "\n\n".join(text_parts) if text_parts else None
    except Exception:
        return None


def extract_text_from_docx(content: bytes) -> Optional[str]:
    """Extract text from Word (.docx) content."""
    try:
        from docx import Document

        doc = Document(io.BytesIO(content))
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_parts.append(cell.text)
        return "\n\n".join(text_parts) if text_parts else None
    except Exception:
        return None


def extract_text_from_xlsx(content: bytes) -> Optional[str]:
    """Extract text from Excel (.xlsx) content."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        text_parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text_parts.append(f"[Sheet: {sheet}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    text_parts.append(" | ".join(cells))
        wb.close()
        return "\n".join(text_parts) if text_parts else None
    except Exception:
        return None


def extract_text_from_csv(content: bytes) -> Optional[str]:
    """Extract text from CSV content."""
    import csv

    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = content.decode("latin-1")
        except Exception:
            return None

    try:
        reader = csv.reader(io.StringIO(text))
        rows = []
        for row in reader:
            if any(cell.strip() for cell in row):
                rows.append(" | ".join(row))
        return "\n".join(rows) if rows else None
    except Exception:
        return None


def extract_text_from_pptx(content: bytes) -> Optional[str]:
    """Extract text from PowerPoint (.pptx) content."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_texts.append(paragraph.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_texts.append(" | ".join(cells))
            if slide_texts:
                text_parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
        return "\n\n".join(text_parts) if text_parts else None
    except Exception:
        return None


def get_openai_client() -> OpenAI:
    """Get OpenAI client instance."""
    return OpenAI(api_key=settings.openai_api_key)


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Split text into overlapping chunks."""
    if len(text) <= chunk_size:
        return [text]

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        start = end - overlap

    return chunks


def generate_embedding(text: str) -> List[float]:
    """Generate embedding for a single text using OpenAI."""
    if not settings.openai_api_key:
        raise ValueError("OpenAI API key not configured")

    client = get_openai_client()
    response = client.embeddings.create(
        model=settings.embedding_model,
        input=text,
    )
    return response.data[0].embedding


def generate_embeddings_batch(texts: List[str]) -> List[List[float]]:
    """Generate embeddings for multiple texts in a single API call."""
    if not settings.openai_api_key:
        raise ValueError("OpenAI API key not configured")

    if not texts:
        return []

    client = get_openai_client()
    response = client.embeddings.create(
        model=settings.embedding_model,
        input=texts,
    )
    return [item.embedding for item in response.data]


async def generate_and_store_embeddings(
    db: Session, file: File, content: bytes, content_type: str
) -> bool:
    """Generate embeddings for file content and store them.

    Returns True if embeddings were generated successfully (or file was empty),
    False if embedding generation failed.
    """
    text = None

    # Handle PDF files
    if content_type == "application/pdf":
        text = extract_text_from_pdf(content)
        if text is None:
            logger.warning("Failed to extract text from PDF %s", file.id)
            return False
    # Handle Word documents (.docx)
    elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        text = extract_text_from_docx(content)
        if text is None:
            logger.warning("Failed to extract text from Word doc %s", file.id)
            return False
    # Handle Excel (.xlsx)
    elif content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        text = extract_text_from_xlsx(content)
        if text is None:
            logger.warning("Failed to extract text from Excel %s", file.id)
            return False
    # Handle PowerPoint (.pptx)
    elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text = extract_text_from_pptx(content)
        if text is None:
            logger.warning("Failed to extract text from PowerPoint %s", file.id)
            return False
    # Handle CSV
    elif content_type == "text/csv":
        text = extract_text_from_csv(content)
        if text is None:
            logger.warning("Failed to extract text from CSV %s", file.id)
            return False
    else:
        # Decode content to text
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            try:
                text = content.decode("latin-1")
            except Exception:
                logger.warning("Failed to decode file %s for embedding", file.id)
                return False

    if not text or not text.strip():
        return True  # Empty file, nothing to embed

    # Chunk the text
    chunks = chunk_text(text)

    # Generate embeddings
    try:
        embeddings = generate_embeddings_batch(chunks)
    except Exception:
        logger.exception("Failed to generate embeddings for file %s", file.id)
        return False

    # Store embeddings
    for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
        file_embedding = FileEmbedding(
            file_id=file.id,
            chunk_index=i,
            chunk_text=chunk,
            embedding=embedding,
        )
        db.add(file_embedding)

    return True


def suggest_folder(filename: str, content_type: str, text_preview: str, existing_folders: List[str]) -> str:
    """Use LLM to suggest a folder for a file based on its content."""
    if not settings.openai_api_key:
        return "/"

    try:
        client = get_openai_client()
        folder_list = "\n".join(existing_folders) if existing_folders else "(none yet)"
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You organize files into folders. Return ONLY the folder path (e.g., /reports/finance/). "
                        "Rules:\n"
                        "- Use existing folders when they fit\n"
                        "- Create new folders only when no existing folder is appropriate\n"
                        "- Keep folder names short and lowercase\n"
                        "- Max 2 levels deep\n"
                        "- Always start with / and end with /\n"
                        "- Return ONLY the path, nothing else"
                    ),
                },
                {
                    "role": "user",
                    "content": (
                        f"Filename: {filename}\n"
                        f"Type: {content_type}\n"
                        f"Content preview: {text_preview[:500]}\n"
                        f"Existing folders:\n{folder_list}"
                    ),
                },
            ],
            max_tokens=50,
            temperature=0,
        )
        folder = response.choices[0].message.content.strip()
        # Validate the response
        if not folder.startswith("/"):
            folder = "/" + folder
        if not folder.endswith("/"):
            folder = folder + "/"
        return folder
    except Exception:
        logger.exception("Auto-organize failed for %s", filename)
        return "/"


def search_embeddings(db: Session, token_id: str, query: str, limit: int = 10) -> List[dict]:
    """Search for files matching the query using vector similarity."""
    # Generate embedding for query
    query_embedding = generate_embedding(query)

    # Search using pgvector cosine distance
    from sqlalchemy import select, func
    from mvp.models import FileEmbedding, File

    # Query with vector similarity
    results = (
        db.query(
            FileEmbedding,
            File,
            FileEmbedding.embedding.cosine_distance(query_embedding).label("distance"),
        )
        .join(File, FileEmbedding.file_id == File.id)
        .filter(File.token_id == token_id)
        .order_by("distance")
        .limit(limit)
        .all()
    )

    # Deduplicate by file and format results
    seen_files = set()
    formatted_results = []

    for embedding, file, distance in results:
        if file.id in seen_files:
            continue
        seen_files.add(file.id)

        formatted_results.append({
            "file_id": str(file.id),
            "filename": file.filename,
            "folder": file.folder,
            "content_type": file.content_type,
            "relevance_score": 1 - distance,  # Convert distance to similarity
            "matched_chunk": embedding.chunk_text[:200] + "..." if len(embedding.chunk_text) > 200 else embedding.chunk_text,
        })

    return formatted_results
