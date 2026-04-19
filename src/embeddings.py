"""Embeddings generation and management using Google Gemini."""

import io
import logging
from typing import List, Optional

from google import genai
from google.genai import types

from src.config import settings
from src.embedding_jobs import EmbeddingJobError, EmbeddingWrite
from src.models import File, FileEmbedding

logger = logging.getLogger(__name__)

# Chunk size for text splitting (in characters)
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
IMAGE_CAPTION_PROMPT = (
    "Describe this image for search indexing in one short sentence. "
    "Include only concrete visible subjects, setting, actions, colors, "
    "and readable text. Do not guess."
)

# Content types that can be embedded directly via Gemini multimodal
IMAGE_CONTENT_TYPES = {
    "image/png",
    "image/jpeg",
    "image/gif",
    "image/webp",
}

VIDEO_CONTENT_TYPES = {
    "video/mp4",
    "video/quicktime",  # .mov
}

AUDIO_CONTENT_TYPES = {
    "audio/mpeg",  # .mp3
    "audio/wav",
    "audio/x-wav",
}

# All types that bypass text extraction and embed raw bytes directly
MULTIMODAL_CONTENT_TYPES = IMAGE_CONTENT_TYPES | VIDEO_CONTENT_TYPES | AUDIO_CONTENT_TYPES


def extract_text_from_pdf(content: bytes) -> Optional[str]:
    """Extract text from PDF content."""
    try:
        import pdfplumber

        with pdfplumber.open(io.BytesIO(content)) as pdf:
            text_parts = []
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            return "\n\n".join(text_parts) if text_parts else None
    except Exception:
        return None


def extract_text_from_docx(content: bytes) -> Optional[str]:
    """Extract text from Word (.docx) content."""
    try:
        from docx import Document

        doc = Document(io.BytesIO(content))
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_parts.append(cell.text)
        return "\n\n".join(text_parts) if text_parts else None
    except Exception:
        return None


def extract_text_from_xlsx(content: bytes) -> Optional[str]:
    """Extract text from Excel (.xlsx) content."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        text_parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text_parts.append(f"[Sheet: {sheet}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    text_parts.append(" | ".join(cells))
        wb.close()
        return "\n".join(text_parts) if text_parts else None
    except Exception:
        return None


def extract_text_from_csv(content: bytes) -> Optional[str]:
    """Extract text from CSV content."""
    import csv

    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = content.decode("latin-1")
        except Exception:
            return None

    try:
        reader = csv.reader(io.StringIO(text))
        rows = []
        for row in reader:
            if any(cell.strip() for cell in row):
                rows.append(" | ".join(row))
        return "\n".join(rows) if rows else None
    except Exception:
        return None


def extract_text_from_pptx(content: bytes) -> Optional[str]:
    """Extract text from PowerPoint (.pptx) content."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_texts.append(paragraph.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_texts.append(" | ".join(cells))
            if slide_texts:
                text_parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
        return "\n\n".join(text_parts) if text_parts else None
    except Exception:
        return None


def get_gemini_client() -> genai.Client:
    """Get Gemini client instance."""
    return genai.Client(api_key=settings.google_api_key)


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Split text into overlapping chunks."""
    if len(text) <= chunk_size:
        return [text]

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        start = end - overlap

    return chunks


def generate_image_caption(content: bytes, content_type: str) -> Optional[str]:
    """Generate a short caption for an image to improve retrieval."""
    if not settings.google_api_key:
        raise ValueError("Google API key not configured")

    client = get_gemini_client()
    result = client.models.generate_content(
        model=settings.image_caption_model,
        contents=[
            IMAGE_CAPTION_PROMPT,
            types.Part.from_bytes(data=content, mime_type=content_type),
        ],
        config=types.GenerateContentConfig(max_output_tokens=500, temperature=0.1),
    )
    caption = (result.text or "").strip()
    return caption or None


def embed_query(text: str) -> List[float]:
    """Generate embedding for a single text query using Gemini."""
    if not settings.google_api_key:
        raise ValueError("Google API key not configured")

    client = get_gemini_client()
    result = client.models.embed_content(
        model=settings.embedding_model,
        contents=f"task: search query | query: {text}",
        config=types.EmbedContentConfig(
            output_dimensionality=settings.embedding_dimensions,
        ),
    )
    return list(result.embeddings[0].values)


def generate_embedding(text: str) -> List[float]:
    """Backward-compatible alias for query embedding."""
    return embed_query(text)


def generate_embeddings_batch(texts: List[str], filename: str) -> List[List[float]]:
    """Generate embeddings for multiple texts."""
    if not settings.google_api_key:
        raise ValueError("Google API key not configured")

    if not texts:
        return []

    client = get_gemini_client()
    embeddings = []
    for text in texts:
        result = client.models.embed_content(
            model=settings.embedding_model,
            contents=f"title: {filename} | text: {text}",
            config=types.EmbedContentConfig(
                output_dimensionality=settings.embedding_dimensions,
            ),
        )
        embeddings.append(list(result.embeddings[0].values))
    return embeddings


def generate_multimodal_embedding(
    content: bytes,
    content_type: str,
    filename: str,
    caption: Optional[str] = None,
) -> List[float]:
    """Generate embedding for an image, video, or audio file using Gemini multimodal embedding."""
    if not settings.google_api_key:
        raise ValueError("Google API key not configured")

    if content_type in IMAGE_CONTENT_TYPES:
        text_context = f"filename: {filename}"
        if caption:
            text_context = f"{text_context} | caption: {caption}"
        embed_content: types.Content | types.Part = types.Content(parts=[
            types.Part.from_text(text=text_context),
            types.Part.from_bytes(data=content, mime_type=content_type),
        ])
    else:
        embed_content = types.Part.from_bytes(data=content, mime_type=content_type)

    client = get_gemini_client()
    result = client.models.embed_content(
        model=settings.embedding_model,
        contents=embed_content,
        config=types.EmbedContentConfig(
            output_dimensionality=settings.embedding_dimensions,
        ),
    )
    return list(result.embeddings[0].values)


def embed_file_content(file: File, content: bytes, content_type: str) -> List[EmbeddingWrite]:
    """Generate the replacement embedding payload for a file."""
    # Handle multimodal files (image/video/audio) — embed directly via Gemini
    if content_type in MULTIMODAL_CONTENT_TYPES:
        caption = None
        if content_type in IMAGE_CONTENT_TYPES:
            try:
                caption = generate_image_caption(content, content_type)
            except Exception:
                logger.exception("Failed to generate image caption for file %s", file.id)

        try:
            embedding = generate_multimodal_embedding(
                content,
                content_type,
                file.filename,
                caption=caption,
            )
        except Exception as exc:
            logger.exception("Failed to generate multimodal embedding for file %s", file.id)
            raise EmbeddingJobError(
                "embedding_provider_error",
                f"Failed to generate multimodal embedding for {file.filename}",
            ) from exc

        # Determine label for the chunk text placeholder
        if content_type in IMAGE_CONTENT_TYPES:
            label = caption or f"[image: {file.filename}]"
        elif content_type in VIDEO_CONTENT_TYPES:
            label = f"[video: {file.filename}]"
        else:
            label = f"[audio: {file.filename}]"

        return [EmbeddingWrite(chunk_index=0, chunk_text=label, embedding=embedding)]

    # Handle text-based files — extract text, chunk, and embed
    text = None

    if content_type == "application/pdf":
        text = extract_text_from_pdf(content)
        if text is None:
            logger.warning("Failed to extract text from PDF %s", file.id)
            raise EmbeddingJobError("text_extraction_failed", f"Failed to extract text from {file.filename}", retryable=False)
    elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        text = extract_text_from_docx(content)
        if text is None:
            logger.warning("Failed to extract text from Word doc %s", file.id)
            raise EmbeddingJobError("text_extraction_failed", f"Failed to extract text from {file.filename}", retryable=False)
    # Handle Excel (.xlsx)
    elif content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        text = extract_text_from_xlsx(content)
        if text is None:
            logger.warning("Failed to extract text from Excel %s", file.id)
            raise EmbeddingJobError("text_extraction_failed", f"Failed to extract text from {file.filename}", retryable=False)
    # Handle PowerPoint (.pptx)
    elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text = extract_text_from_pptx(content)
        if text is None:
            logger.warning("Failed to extract text from PowerPoint %s", file.id)
            raise EmbeddingJobError("text_extraction_failed", f"Failed to extract text from {file.filename}", retryable=False)
    # Handle CSV
    elif content_type == "text/csv":
        text = extract_text_from_csv(content)
        if text is None:
            logger.warning("Failed to extract text from CSV %s", file.id)
            raise EmbeddingJobError("text_extraction_failed", f"Failed to extract text from {file.filename}", retryable=False)
    elif (
        content_type.startswith("text/")
        or content_type in {"application/json", "application/xml", "text/markdown"}
    ):
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            try:
                text = content.decode("latin-1")
            except Exception as exc:
                logger.warning("Failed to decode file %s for embedding", file.id)
                raise EmbeddingJobError(
                    "text_extraction_failed",
                    f"Failed to decode {file.filename} for embedding",
                    retryable=False,
                ) from exc
    else:
        raise EmbeddingJobError(
            "unsupported_content_type",
            f"Unsupported content type for embedding: {content_type}",
            retryable=False,
        )

    if not text or not text.strip():
        return []

    # Chunk the text
    chunks = chunk_text(text)

    # Generate embeddings
    try:
        embeddings = generate_embeddings_batch(chunks, file.filename)
    except Exception as exc:
        logger.exception("Failed to generate embeddings for file %s", file.id)
        raise EmbeddingJobError(
            "embedding_provider_error",
            f"Failed to generate embeddings for {file.filename}",
        ) from exc

    return [
        EmbeddingWrite(
            chunk_index=i,
            chunk_text=chunk,
            embedding=embedding,
        )
        for i, (chunk, embedding) in enumerate(zip(chunks, embeddings))
    ]


def search_embeddings(db, token_id: str, query: str, limit: int = 10) -> List[dict]:
    """Search for files matching the query using vector similarity."""
    # Generate embedding for query
    query_embedding = embed_query(query)

    # Query with vector similarity
    results = (
        db.query(
            FileEmbedding,
            File,
            FileEmbedding.embedding.cosine_distance(query_embedding).label("distance"),
        )
        .join(File, FileEmbedding.file_id == File.id)
        .filter(File.token_id == token_id)
        .order_by("distance")
        .limit(limit)
        .all()
    )

    # Deduplicate by file and format results
    seen_files = set()
    formatted_results = []

    for embedding, file, distance in results:
        if file.id in seen_files:
            continue
        seen_files.add(file.id)

        formatted_results.append({
            "file_id": str(file.id),
            "filename": file.filename,
            "folder": file.folder,
            "content_type": file.content_type,
            "relevance_score": 1 - distance,  # Convert distance to similarity
            "matched_chunk": embedding.chunk_text[:200] + "..." if len(embedding.chunk_text) > 200 else embedding.chunk_text,
        })

    return formatted_results
